from pptx import Presentation
import six
import copy

from colorama import init
from termcolor import colored
#import numpy as np
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from pptx.dml.color import RGBColor
# new_temp has the template that you want your slides to be formatted to


# old_ppt = Presentation('GPIO.pptx')
################ construct name mappings for two layouts
#https://stackoverflow.com/questions/37340049/how-do-i-print-colored-output-to-the-terminal-in-python/37340245

from collections import namedtuple
DimensionOfShape = namedtuple("DimensionOfShape", "width height")


# use Colorama to make Termcolor work on Windows too


#defining values for slides that may not be imported
left = Inches(1)
top =   Inches(2)
width =  Inches(7.5)
height =  Inches(1.5)

old_layout_names=[]
new_layout_names=[]
slidesDifferentShapes = []
slidesNotAbleToImport = []
slidesNotBeingMapped = []
slidesShapesDimensionDiverged = []
slidesContainLinks = []
slidesContainBackground = []
layout_map={}
new_temp = None
old_ppt = None

init()

#Functions section--------------------------------------------------------------
def corresponding_slide_layout(slide, new_prs):
    legacy_layout_name = slide.slide_layout.name
    new_name = layout_map.get(legacy_layout_name)
    if new_name is None:
        # --- there's a gap in the mapping ---
        return None
    return new_prs.slide_layouts.get_by_name(new_name)

def addErrorSlide(slide):
  add_warning = slide.shapes.add_textbox(left,top,width,height)
  add_warning.rotation = 20
  tf = add_warning.text_frame
  p = tf.paragraphs[0]

  run = p.add_run()
  run.text = "ADD MANUALLY THIS SLIDE"
  font = run.font
  font.size = Pt(72)
  p.aligment = PP_ALIGN.CENTER
  #red color 0xff0000
  font.color.rgb = RGBColor(0XFF,0X00,0X00)
#-------------------------------------------------------------------------------



#This func will get the old presentation and the new presentation template
#and recorded all the layout in both templates and mapped with the respective
#newer layout, at the end of the execution it will return a presentation
def mappingPre(prsOldP, prsNewP):
    #refering to the global variables
    global layout_map
    global old_layout_names
    global new_layout_names
    global slidesDifferentShapes
    global slidesNotAbleToImport
    global slidesNotBeingMapped
    global slidesShapesDimensionDiverged
    global slidesContainLinks
    global slidesContainBackground
    global new_temp
    global old_ppt

    #restarting all global variables
    old_layout_names=[]
    new_layout_names=[]
    slidesDifferentShapes = []
    slidesNotAbleToImport = []
    slidesNotBeingMapped = []
    slidesShapesDimensionDiverged = []
    slidesContainLinks = []
    slidesContainBackground = []
    layout_map={}
    new_temp = None
    old_ppt = None

    new_temp = prsNewP
    old_ppt = prsOldP

    for i, slide_layout in enumerate(old_ppt.slide_masters[0].slide_layouts):
        old_layout_names.append(slide_layout.name)
        print("old layout name {:10} -- {}".format(i, slide_layout.name))

    for i, slide_layout in enumerate(new_temp.slide_masters[0].slide_layouts):
        new_layout_names.append(slide_layout.name)
        print("new layout name {:10} -- {}".format(i, slide_layout.name))

    layout_map=dict(zip(old_layout_names,new_layout_names))
    for key, value in layout_map.items():
        print(key, ' : ', value)

    # print out layout information for all old slides
    print("layout of each old page: \n")
    for i, slide in enumerate(old_ppt.slides):
        print(i,slide.slide_layout.name)

    for slide in old_ppt.slides:
        #old_ppt.slides.index(slide)+1
        dimNewShapes = []
        dimOldShapes = []
        slideNumber = str(old_ppt.slides.index(slide) + 1)
        isThereErrorWithSlide = False
        layout = corresponding_slide_layout(slide, new_temp)
        if layout is None:
            layout = new_temp.slide_layouts[2]  # --- or whatever default --- this might need change
            slidesNotBeingMapped.append(slideNumber)
        new_slide = new_temp.slides.add_slide(layout)
        for shape in new_slide.shapes:
            if(shape.has_text_frame):
                shape.text = "" + "\n"
            #    if(shape)
            dimNewShapes.append(DimensionOfShape(width = round(shape.width.cm), height = round(shape.height.cm)))
        countShapesPerNewerSlide = len(new_slide.shapes)
        countShapesOldSlide = len(slide.shapes)
        if slide.follow_master_background == True:
            for shp in slide.shapes:
                dimOldShapes.append(DimensionOfShape(width = round(shp.width.cm), height = round(shp.height.cm)))
                el = shp.element
                elXML = el.xml
                if "tags" in elXML:
                  slidesNotAbleToImport.append(slideNumber)
                  isThereErrorWithSlide = True
                  addErrorSlide(new_slide)
                  break
                elif "hlinkClick" in elXML:
                  slidesContainLinks.append(slideNumber)
                  isThereErrorWithSlide = True
                  addErrorSlide(new_slide)
                  break
                #adding an old shape shape from the older slide to the newer
                newel = copy.deepcopy(el)
                new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
        else:
            isThereErrorWithSlide = True
            slidesContainBackground.append(slideNumber)
            addErrorSlide(new_slide)

        if(isThereErrorWithSlide == False):
            for key, value in six.iteritems(slide.part.rels):
                # Make sure we don't copy a notesSlide relation as that won't exist
                # Make sure we don't copy over the slideLayout relation as we already have that with the new slide layouts !
                if not "notesSlide" in value.reltype:
                    if not "slideLayout" in value.reltype:
                        new_slide.part.rels.add_relationship(value.reltype, value._target, value.rId)
        if(countShapesPerNewerSlide != countShapesOldSlide ):
          slidesDifferentShapes.append(slideNumber)
        else:
          while(len(dimNewShapes) > 0):
            dimensionMatched = False
            dimNew  = dimNewShapes.pop()
            temp = dimOldShapes
            for dimOld in temp:
              if(dimOld.width == dimNew.width and dimOld.height == dimNew.height):
                dimOldShapes.remove(dimOld)
                dimensionMatched = True
                break
            if(dimensionMatched == False):
              slidesShapesDimensionDiverged.append(slideNumber)
              break
        #trying to add notes
        if(slide.has_notes_slide):
          notes_old_slide = slide.notes_slide
          text_old_frame = notes_old_slide.notes_text_frame

          notes_new_slide = new_slide.notes_slide
          text_new_frame = notes_new_slide.notes_text_frame
          text_new_frame.text = text_old_frame.text

    #Print any warning or error that has been encountered with the presentation
    if(len(slidesDifferentShapes) > 0):
        print(colored("\nWARNING:" ,"yellow") + "The following slides contain more shapes in the older version that might affect the layout of the newer layout:", end = " ")
        for i in slidesDifferentShapes:
            print(i ,end = " ")

    if(len(slidesNotBeingMapped) > 0):
        print(colored("\nERROR" ,"red") + "The following slides used the default layout:", end = " ")
        errorsDetected = True
        for i in slidesNotBeingMapped:
            print(i ,end = " ")

    if(len(slidesShapesDimensionDiverged) > 0):
        print(colored("\nWARNING:" ,"yellow") + "The following slide's shapes have different width or height in comparison to the newer shapes, in slides:", end = " ")
        for i in slidesShapesDimensionDiverged:
            print(i ,end = " ")

    errorsDetected = False
    if(len(slidesNotAbleToImport) > 0):
        errorsDetected = True
        print(colored('\nERROR: ', 'red') +"The following slides could not be imported to the newer layout, please add them manually", end = " ")
        for i in slidesNotAbleToImport:
            print(i, end = " ")

    if(len(slidesContainLinks) > 0):
        errorsDetected = True
        print(colored('\nERROR: ', 'red') +"The following slides could not be imported because they contain links", end = " ")
        for i in slidesContainLinks:
            print(i, end = " ")

    if(len(slidesContainBackground) > 0):
        errorsDetected = True
        print(colored('\nERROR: ', 'red') +"The following slides could not be imported because of a local background in the slides", end = " ")
        for i in slidesContainBackground:
            print(i, end = " ")
    print(colored('\n***************************************************************************************************************************', 'green'))
    # if(errorsDetected == True):
    #   print(colored('There was a problem with your mapping', 'red'))
    #   input("Before continuing ,add the missing slides manually and press Enter to resume this program\n")
    return [new_temp,errorsDetected]
