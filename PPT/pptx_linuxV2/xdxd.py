from pptx import Presentation
import six
import copy

from colorama import init
from termcolor import colored
#import numpy as np
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from pptx.dml.color import RGBColor

import win32com.client
import time
# new_temp has the template that you want your slides to be formatted to
# new_temp = Presentation('new template partial5.pptx')

from collections import namedtuple
DimensionOfShape = namedtuple("DimensionOfShape", "width height")
# old_ppt is the existing ppt file
#old_ppt = Presentation('gaming9v2.pptx')
# filename = 'gaming9v2.pptx'
# old_ppt = Presentation(filename)
filename = None
new_temp = None
old_ppt = None


def RGB(red, green, blue):
    assert 0 <= red <=255
    assert 0 <= green <=255
    assert 0 <= blue <=255
    return red + (green << 8) + (blue << 16)
# use Colorama to make Termcolor work on Windows too
init()


# old_ppt = Presentation('GPIO.pptx')
################ construct name mappings for two layouts
#https://stackoverflow.com/questions/37340049/how-do-i-print-colored-output-to-the-terminal-in-python/37340245
def mappingPreV2(dirPrsOldP, dirPrsNewP,nameModification,directoryToStore):
    #global filename
    # global new_temp
    # global old_ppt

    # #filename = None
    # new_temp = None
    # old_ppt = None
    #
    # #filename = prsOldP
    # new_temp = prsNewP
    # old_ppt = filenameSel


    #new_temp.save('{}_redwood.pptx'.format(filename))

    #
    # dest_path = r'C:\Users\marmor05\Desktop\forkRep\edusuite\PPT\{}_redwood.pptx'.format(filename)
    # source_path = r'C:\Users\marmor05\Desktop\forkRep\edusuite\PPT\{}'.format(filename)


    dest_path = dirPrsNewP
    source_path = dirPrsOldP

    ppt_instance = win32com.client.Dispatch('PowerPoint.Application')
    ppt_instance.Visible = True

    time.sleep(1)
    prs = ppt_instance.Presentations.open(source_path)
    time.sleep(1)
    prs2 = ppt_instance.Presentations.open(dest_path)
    #prs2 = ppt_instance.Presentations.open(dest_path)
    slide_count = len(prs.Slides)

    #prs.Slides.Item(2).Copy()
    #prs2.Slides.Item(2).Select()
    # prs.Slides[1].Select()
    # prs.Slides[1].Copy()

    prs2.Slides.InsertFromFile(source_path,0,1,slide_count)
    for i in range(slide_count):
        if(prs.Slides[i].FollowMasterBackground == False):
          print("the slide back" + str(i + 1))
          #remove the if statement , this is just for GG gaming course
          if(i == slide_count - 1):
             prs2.Slides[i].CustomLayout = prs2.Designs[i].SlideMaster.CustomLayouts[len(prs2.Designs[i].SlideMaster.CustomLayouts)-2]
          else:
              prs2.Slides[i].FollowMasterBackground = False
              for shape in prs.Slides[i].Shapes:
                shape.Visible = False
              #prs.Slides[i].CustomLayout = prs.Designs[i].SlideMaster.CustomLayouts[1]
              prs.Slides[i].DisplayMasterShapes = False
              prs.Slides[i].Export(directoryToStore + "/slide.png","PNG",13333,7500)
              #Disable the followmaster of the destation pres to allow insertation of background
              prs2.Slides[i].FollowMasterBackground = False
              #prs2.Slides[i].Background.Fill.UserPicture(r"C:\Users\marmor05\Desktop\labs\bookCover.png")
              prs2.Slides[i].Background.Fill.UserPicture(directoryToStore + "/slide.png")

              #prs2.Slides[i].Fill.Solid
              #prs2.Slides[i].Background.Fill.ForeColor.RGB = RGB(255, 0, 0)
              #prs2.Slides[i].Background.Fill.PresetTextured("msoShapeCan",90,90,40,80)
              #print(str(prs.Slides[i].Background.Fill.PresetGradient))
              #print(str(prs2.Slides[i].Background))

    #prs2.SaveAs(r'C:\Users\marmor05\Desktop\forkRep\edusuite\PPT\{}_redwood.pptx'.format(filename))
    prs2.SaveAs(nameModification)
    prs.Close()
    prs2.Close()
    ppt_instance.Quit()
    del ppt_instance
    # return [prs2,False]

# print out layout information fo
